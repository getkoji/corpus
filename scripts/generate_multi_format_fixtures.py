#!/usr/bin/env python3
"""Generate synthetic docx/xlsx/pptx fixtures for the multi_format corpus.

Each fixture contains the same deterministic content (an invoice from
Meridian Supply Co. to Horizon Holdings LLC) so expected JSONs can
assert identical values across formats. The point of the category is
to verify that `koji parse` produces extractable markdown from each
source format, not to test format-specific content extraction.

Dependencies: python-docx, openpyxl, python-pptx. Run via uv:

    uv run --with python-docx --with openpyxl --with python-pptx \\
        python scripts/generate_multi_format_fixtures.py
"""
from __future__ import annotations

from datetime import date
from pathlib import Path

CORPUS_ROOT = Path(__file__).resolve().parent.parent
SOURCES_DIR = CORPUS_ROOT / "multi_format" / "sources"

# The deterministic invoice payload every fixture shares. Expected JSONs
# assert exactly these values regardless of source format.
INVOICE = {
    "merchant_name": "Meridian Supply Co.",
    "date": "2026-03-22",
    "date_display": "March 22, 2026",
    "bill_to": "Horizon Holdings LLC",
    "invoice_number": "INV-2026-04412",
    "line_items": [
        ("Office furniture — desks", 12, 450.00, 5400.00),
        ("Office furniture — chairs", 24, 180.00, 4320.00),
        ("Delivery and assembly", 1, 800.00, 800.00),
    ],
    "subtotal": 10520.00,
    "tax_rate": 0.0825,
    "tax": 867.90,
    "total_amount": 11387.90,
    "currency": "USD",
}


def build_docx_invoice(out_path: Path) -> None:
    """A straightforward Word invoice — heading, address block, table, totals."""
    from docx import Document
    from docx.shared import Pt

    doc = Document()

    title = doc.add_heading("INVOICE", level=0)
    title.alignment = 0

    merchant = doc.add_paragraph()
    merchant_run = merchant.add_run(INVOICE["merchant_name"])
    merchant_run.bold = True
    merchant_run.font.size = Pt(14)
    doc.add_paragraph("2200 Commerce Avenue")
    doc.add_paragraph("Dallas, TX 75201")
    doc.add_paragraph("Phone: (214) 555-0188")

    doc.add_paragraph()
    doc.add_paragraph(f"Invoice #: {INVOICE['invoice_number']}")
    doc.add_paragraph(f"Invoice Date: {INVOICE['date_display']}")
    doc.add_paragraph("Due Date: April 21, 2026")

    doc.add_paragraph()
    bill = doc.add_paragraph()
    bill_run = bill.add_run("Bill To:")
    bill_run.bold = True
    doc.add_paragraph(INVOICE["bill_to"])
    doc.add_paragraph("410 Corporate Drive")
    doc.add_paragraph("Plano, TX 75024")

    doc.add_paragraph()
    table = doc.add_table(rows=1, cols=4)
    table.style = "Light Grid Accent 1"
    header_row = table.rows[0].cells
    header_row[0].text = "Item"
    header_row[1].text = "Qty"
    header_row[2].text = "Unit Price"
    header_row[3].text = "Total"
    for desc, qty, unit, total in INVOICE["line_items"]:
        row = table.add_row().cells
        row[0].text = desc
        row[1].text = str(qty)
        row[2].text = f"${unit:,.2f}"
        row[3].text = f"${total:,.2f}"

    doc.add_paragraph()
    doc.add_paragraph(f"Subtotal: ${INVOICE['subtotal']:,.2f}")
    doc.add_paragraph(f"Tax ({INVOICE['tax_rate'] * 100:.2f}%): ${INVOICE['tax']:,.2f}")
    total_p = doc.add_paragraph()
    total_run = total_p.add_run(f"Total Due: ${INVOICE['total_amount']:,.2f}")
    total_run.bold = True

    doc.add_paragraph()
    doc.add_paragraph(
        "Payment terms: Net 30. Please remit to Meridian Supply Co. at the address above. "
        "Thank you for your business."
    )

    doc.save(out_path)


def build_xlsx_invoice(out_path: Path) -> None:
    """Excel invoice — a header block in rows 1-10, then a line-item table."""
    from openpyxl import Workbook
    from openpyxl.styles import Font

    wb = Workbook()
    ws = wb.active
    ws.title = "Invoice"

    ws["A1"] = "INVOICE"
    ws["A1"].font = Font(size=18, bold=True)

    ws["A3"] = INVOICE["merchant_name"]
    ws["A3"].font = Font(bold=True)
    ws["A4"] = "2200 Commerce Avenue"
    ws["A5"] = "Dallas, TX 75201"
    ws["A6"] = "Phone: (214) 555-0188"

    ws["A8"] = "Invoice #:"
    ws["B8"] = INVOICE["invoice_number"]
    ws["A9"] = "Invoice Date:"
    ws["B9"] = INVOICE["date_display"]
    ws["A10"] = "Due Date:"
    ws["B10"] = "April 21, 2026"

    ws["A12"] = "Bill To:"
    ws["A12"].font = Font(bold=True)
    ws["A13"] = INVOICE["bill_to"]
    ws["A14"] = "410 Corporate Drive"
    ws["A15"] = "Plano, TX 75024"

    header_row = 17
    ws.cell(row=header_row, column=1, value="Item").font = Font(bold=True)
    ws.cell(row=header_row, column=2, value="Qty").font = Font(bold=True)
    ws.cell(row=header_row, column=3, value="Unit Price").font = Font(bold=True)
    ws.cell(row=header_row, column=4, value="Total").font = Font(bold=True)
    for i, (desc, qty, unit, total) in enumerate(INVOICE["line_items"], start=1):
        r = header_row + i
        ws.cell(row=r, column=1, value=desc)
        ws.cell(row=r, column=2, value=qty)
        ws.cell(row=r, column=3, value=unit)
        ws.cell(row=r, column=4, value=total)

    totals_row = header_row + len(INVOICE["line_items"]) + 2
    ws.cell(row=totals_row, column=3, value="Subtotal:").font = Font(bold=True)
    ws.cell(row=totals_row, column=4, value=INVOICE["subtotal"])
    ws.cell(row=totals_row + 1, column=3, value=f"Tax ({INVOICE['tax_rate'] * 100:.2f}%):").font = Font(bold=True)
    ws.cell(row=totals_row + 1, column=4, value=INVOICE["tax"])
    ws.cell(row=totals_row + 2, column=3, value="Total Due:").font = Font(bold=True)
    ws.cell(row=totals_row + 2, column=4, value=INVOICE["total_amount"]).font = Font(bold=True)

    # Widen columns a bit so the output looks reasonable when opened.
    for col, width in (("A", 32), ("B", 16), ("C", 16), ("D", 16)):
        ws.column_dimensions[col].width = width

    wb.save(out_path)


def build_pptx_invoice(out_path: Path) -> None:
    """PowerPoint invoice across three slides: cover, line items, totals."""
    from pptx import Presentation
    from pptx.util import Inches, Pt

    prs = Presentation()
    blank = prs.slide_layouts[6]

    def add_textbox(slide, left, top, width, height, text, bold=False, size=18):
        box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        return box

    # Slide 1 — cover
    s1 = prs.slides.add_slide(blank)
    add_textbox(s1, 0.5, 0.4, 9, 1.0, "INVOICE", bold=True, size=36)
    add_textbox(s1, 0.5, 1.4, 9, 0.5, INVOICE["merchant_name"], bold=True, size=22)
    add_textbox(s1, 0.5, 2.0, 9, 0.5, "2200 Commerce Avenue, Dallas, TX 75201", size=16)
    add_textbox(s1, 0.5, 2.5, 9, 0.5, "Phone: (214) 555-0188", size=16)
    add_textbox(s1, 0.5, 3.4, 9, 0.5, f"Invoice #: {INVOICE['invoice_number']}", size=16)
    add_textbox(s1, 0.5, 3.9, 9, 0.5, f"Invoice Date: {INVOICE['date_display']}", size=16)
    add_textbox(s1, 0.5, 4.4, 9, 0.5, "Due Date: April 21, 2026", size=16)
    add_textbox(s1, 0.5, 5.2, 9, 0.5, "Bill To:", bold=True, size=18)
    add_textbox(s1, 0.5, 5.7, 9, 0.5, INVOICE["bill_to"], size=16)
    add_textbox(s1, 0.5, 6.2, 9, 0.5, "410 Corporate Drive, Plano, TX 75024", size=14)

    # Slide 2 — line items as a shape-based table
    s2 = prs.slides.add_slide(blank)
    add_textbox(s2, 0.5, 0.4, 9, 0.8, "Line Items", bold=True, size=28)
    rows = len(INVOICE["line_items"]) + 1
    cols = 4
    table_shape = s2.shapes.add_table(rows, cols, Inches(0.5), Inches(1.4), Inches(9.0), Inches(0.6 * rows))
    table = table_shape.table
    headers = ["Item", "Qty", "Unit Price", "Total"]
    for c, header in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = header
    for r, (desc, qty, unit, total) in enumerate(INVOICE["line_items"], start=1):
        table.cell(r, 0).text = desc
        table.cell(r, 1).text = str(qty)
        table.cell(r, 2).text = f"${unit:,.2f}"
        table.cell(r, 3).text = f"${total:,.2f}"

    # Slide 3 — totals
    s3 = prs.slides.add_slide(blank)
    add_textbox(s3, 0.5, 0.4, 9, 0.8, "Totals", bold=True, size=28)
    add_textbox(s3, 0.5, 1.6, 9, 0.5, f"Subtotal: ${INVOICE['subtotal']:,.2f}", size=22)
    add_textbox(s3, 0.5, 2.3, 9, 0.5, f"Tax ({INVOICE['tax_rate'] * 100:.2f}%): ${INVOICE['tax']:,.2f}", size=22)
    add_textbox(s3, 0.5, 3.0, 9, 0.5, f"Total Due: ${INVOICE['total_amount']:,.2f}", bold=True, size=26)
    add_textbox(s3, 0.5, 4.2, 9, 2.0,
                "Payment terms: Net 30. Please remit to Meridian Supply Co. "
                "at the address above. Thank you for your business.",
                size=14)

    prs.save(out_path)


def main() -> int:
    SOURCES_DIR.mkdir(parents=True, exist_ok=True)

    build_docx_invoice(SOURCES_DIR / "meridian_invoice.docx")
    print(f"wrote {SOURCES_DIR / 'meridian_invoice.docx'}")

    build_xlsx_invoice(SOURCES_DIR / "meridian_invoice.xlsx")
    print(f"wrote {SOURCES_DIR / 'meridian_invoice.xlsx'}")

    build_pptx_invoice(SOURCES_DIR / "meridian_invoice.pptx")
    print(f"wrote {SOURCES_DIR / 'meridian_invoice.pptx'}")

    return 0


if __name__ == "__main__":
    raise SystemExit(main())
